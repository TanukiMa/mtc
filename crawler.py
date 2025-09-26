# crawler.py
import re
import os
import io
import requests
from datetime import datetime, timedelta, timezone
from urllib.parse import urljoin, urlparse
from concurrent.futures import ThreadPoolExecutor

# --- 外部ライブラリのインポート ---
from bs4 import BeautifulSoup
import pdfplumber
import docx
from pptx import Presentation
from supabase import create_client, Client
from sudachipy import tokenizer
from sudachipy import dictionary

# --- 設定項目 ---
START_URL = "https://www.mhlw.go.jp/"
TARGET_DOMAIN = "www.mhlw.go.jp"
MAX_URLS_TO_CRAWL = 2000000000000
MAX_WORKERS = 8
RECRAWL_DAYS = 0
REQUEST_TIMEOUT = 15

# --- Sudachiの初期化 ---
print("Sudachi tokenizerを初期化しています...")
tokenizer_obj = dictionary.Dictionary(dict_type="full").create(mode=tokenizer.Tokenizer.SplitMode.C)
print("初期化が完了しました。")


# --- テキストクリーニング関数 ---
def clean_text(text: str) -> str:
    """テキストから不要な改行やスペースを削除・整形する"""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text).strip()
    return text


# --- テキスト抽出関数 ---
def get_text(content: bytes, content_type: str) -> str:
    """コンテンツの種別に応じてテキストを抽出する"""
    text = ""
    if "html" in content_type:
        soup = BeautifulSoup(content, 'html.parser')
        for s in soup(['script', 'style']): s.decompose()
        text = ' '.join(soup.stripped_strings)
    elif "pdf" in content_type:
        try:
            with io.BytesIO(content) as pdf_file:
                with pdfplumber.open(pdf_file) as pdf:
                    all_text = []
                    for page in pdf.pages:
                        extracted = page.extract_text()
                        if extracted:
                            all_text.append(extracted)
                    text = "\n".join(all_text)
        except Exception as e:
            print(f"  [!] PDF解析エラー (pdfplumber): {e}")
    elif "vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type:  # docx
        try:
            doc = docx.Document(io.BytesIO(content))
            all_text = [para.text for para in doc.paragraphs]
            text = "\n".join(all_text)
        except Exception as e:
            print(f"  [!] DOCX解析エラー: {e}")
    elif "vnd.openxmlformats-officedocument.presentationml.presentation" in content_type:  # pptx
        try:
            prs = Presentation(io.BytesIO(content))
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        all_text.append(shape.text)
            text = "\n".join(all_text)
        except Exception as e:
            print(f"  [!] PPTX解析エラー: {e}")
    return clean_text(text)


# --- Sudachi連携関数 ---
def analyze_with_sudachi(text: str) -> list:
    """sudachipyを使い、未知の普通名詞を抽出する。長文は分割して処理する。"""
    if not text.strip():
        return []

    chunk_size = 40000  # sudachipyのバイト制限(約49KB)より安全に小さい値
    words = []

    try:
        # テキストをチャンクに分割してループ処理
        for i in range(0, len(text), chunk_size):
            chunk = text[i:i + chunk_size]
            
            # 各チャンクを形態素解析
            morphemes = tokenizer_obj.tokenize(chunk)
            
            for m in morphemes:
                pos_info = m.part_of_speech()
                # 未知語(OOV) かつ 品詞が「名詞,普通名詞」のものを抽出
                if m.is_oov() and pos_info[0] == "名詞" and pos_info[1] == "普通名詞":
                    # 専門用語として意味をなしやすいよう、2文字以上の単語に限定
                    if len(m.surface()) > 1:
                        words.append({
                            "word": m.surface(),
                            "pos": ",".join(pos_info[0:4])
                        })
    except Exception as e:
        # sudachipyの内部エラーを捕捉
        print(f"  [!] Sudachi解析エラー: {e}")

    return words


# --- クローラー本体 ---
def process_url(url: str, supabase_url: str, supabase_key: str, stop_words_set: set):
    """単一のURLを処理するワーカー関数"""
    supabase: Client = create_client(supabase_url, supabase_key)

    try:
        print(f"[*] 処理中: {url}")
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'}
        response = requests.get(url, timeout=REQUEST_TIMEOUT, headers=headers)
        response.raise_for_status()
        content_type = response.headers.get("content-type", "").lower()

        text = get_text(response.content, content_type)
        if not text:
            print(f"  [-] テキスト抽出スキップ: {url}")
            return []

        new_words = analyze_with_sudachi(text)
        
        if stop_words_set:
            filtered_words = [
                word_data for word_data in new_words 
                if word_data["word"] not in stop_words_set
            ]
        else:
            filtered_words = new_words

        if filtered_words:
            print(f"  [+] 発見した未知語候補 ({len(filtered_words)}件): {[w['word'] for w in filtered_words]}")
            
            for word_data in filtered_words:
                try:
                    upsert_response = supabase.table("unique_words").upsert(
                        {"word": word_data["word"], "pos": word_data["pos"]},
                        on_conflict="word"
                    ).execute()

                    if upsert_response.data:
                        word_id = upsert_response.data[0]['id']
                        supabase.table("word_occurrences").insert({
                            "word_id": word_id,
                            "source_url": url
                        }).execute()
                except Exception as e:
                    if "duplicate key value violates unique constraint" not in str(e):
                        print(f"  [!] DB保存エラー(insert): {e}")

        found_links = []
        if "html" in content_type:
            soup = BeautifulSoup(response.content, 'html.parser')
            for a_tag in soup.find_all('a', href=True):
                link = urljoin(url, a_tag['href']).split('#')[0]
                if urlparse(link).netloc == TARGET_DOMAIN:
                    found_links.append(link)
        return found_links
    except requests.RequestException as e:
        print(f"  [!] HTTPエラー: {url} - {e}")
    except Exception as e:
        print(f"  [!] 不明なエラー: {url} - {e}")
    return []


# --- メイン処理 ---
def main():
    """スクリプトのメイン処理"""
    supabase_url: str = os.environ.get("SUPABASE_URL")
    supabase_key: str = os.environ.get("SUPABASE_KEY")
    if not supabase_url or not supabase_key:
        raise ValueError("環境変数 SUPABASE_URL と SUPABASE_KEY を設定してください。")

    supabase_main: Client = create_client(supabase_url, supabase_key)

    print("--- クロール処理開始 ---")

    print("[*] 除外リスト(stop_words)を読み込んでいます...")
    try:
        response = supabase_main.table("stop_words").select("word").execute()
        stop_words_set = {item['word'] for item in response.data}
        print(f"  [+] {len(stop_words_set)}件の除外ワードを読み込みました。")
    except Exception as e:
        print(f"  [!] 除外リストの読み込み中にエラー: {e}")
        stop_words_set = set()

    print(f"[*] {RECRAWL_DAYS}日以上前の処理済みURLを削除します...")
    threshold_date = datetime.now(timezone.utc) - timedelta(days=RECRAWL_DAYS)
    try:
        supabase_main.table("crawled_urls").delete().lt("crawled_at", threshold_date.isoformat()).execute()
    except Exception as e:
        print(f"  [!] 古いURLの削除中にエラーが発生しました: {e}")

    try:
        response = supabase_main.table("crawled_urls").select("url", count='exact').execute()
        crawled_urls = {item['url'] for item in response.data}
        print(f"[*] 現在 {len(crawled_urls)} 件のURLが処理済みです (重複除く)。")
    except Exception as e:
        print(f"  [!] 処理済みURLの取得中にエラーが発生しました: {e}")
        crawled_urls = set()

    urls_to_crawl = {START_URL}
    processed_count = 0
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        while urls_to_crawl and processed_count < MAX_URLS_TO_CRAWL:
            urls_now = list(urls_to_crawl - crawled_urls)
            urls_to_crawl.clear()
            if not urls_now:
                break

            future_to_url = {
                executor.submit(process_url, url, supabase_url, supabase_key, stop_words_set): url
                for url in urls_now
            }

            for future in future_to_url:
                url = future_to_url[future]
                try:
                    new_links = future.result()
                    if new_links:
                        urls_to_crawl.update(new_links)
                    crawled_urls.add(url)
                    supabase_main.table("crawled_urls").insert({"url": url}).execute()
                    processed_count += 1
                    if processed_count >= MAX_URLS_TO_CRAWL:
                        print("[*] 処理上限数に達しました。")
                        break
                except Exception as e:
                    print(f"[!] future.result()でエラー: {url} - {e}")

    print(f"\n--- クロール処理終了 ---")
    print(f"今回処理したURL数: {processed_count}")


if __name__ == "__main__":
    main()
